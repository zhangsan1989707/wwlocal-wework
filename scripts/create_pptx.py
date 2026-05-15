#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_shape.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. 封面
    add_title_slide(prs, "政务微信数据审计平台", "系统功能演示与产品介绍\n\n项目版本：v1.1.0")

    # 2. 目录
    add_content_slide(prs, "目录", [
        "1. 系统概述与核心价值",
        "2. 系统架构",
        "3. 功能模块介绍",
        "4. 核心功能演示",
        "5. 安全与运维保障",
        "6. 部署与使用"
    ])

    # 3. 系统概述
    add_content_slide(prs, "系统概述 - 解决什么问题", [
        "政务微信开放接口返回的业务数据经过 RSA+AES 双层加密",
        "需要一个平台统一解密、存储、查询和审计这些数据",
        "数据不出内网，所有组件部署在本地服务器",
        "支持多维度查询、组织架构关联、操作审计"
    ])

    # 4. 核心能力
    add_content_slide(prs, "系统概述 - 核心能力", [
        "自动从政务微信 API 拉取加密日志",
        "RSA+AES 双层解密后存储到本地 MySQL",
        "多维度查询、筛选、导出（CSV）",
        "组织架构同步与人员关联",
        "操作审计日志自动记录",
        "密钥热切换，支持多版本管理"
    ])

    # 5. 系统架构
    add_content_slide(prs, "系统架构", [
        "技术栈：Go 1.21 + Vue 3 + Element Plus + MySQL 8.0 + Docker",
        "三容器编排：MySQL + 后端 (Echo) + 前端 (Nginx)",
        "端口映射：MySQL 3307，后端 3010，前端 5173",
        "Nginx 反向代理 /api/ 路径到后端",
        "支持 Vue SPA 路由，proxy_read_timeout 300s 适应长同步"
    ])

    # 6. 后端分层
    add_content_slide(prs, "后端分层架构", [
        "Handler 层：HTTP 处理，定义本地 Request struct",
        "Service 层：业务逻辑编排，核心解密管道",
        "Repository 层：数据访问，动态建表和去重",
        "Model 层：数据模型定义",
        "Crypto 层：RSA PKCS1v15 + AES-128-CBC 解密实现",
        "中间件：JWT 认证、操作审计日志"
    ])

    # 7. 前端架构
    add_content_slide(prs, "前端架构", [
        "Vue 3.5 + Composition API + TypeScript",
        "Element Plus 组件库，自动导入",
        "无路由库，通过 v-if 切换视图",
        "无全局状态管理，组件内独立状态",
        "Axios 封装，baseURL=/api/v1，自动附加 JWT",
        "Vite 8 开发服务器，dev 时代理 /api"
    ])

    # 8. 功能总览
    add_content_slide(prs, "功能总览", [
        "总览看板：KPI 概览、同步状态、异常告警、部门活跃度",
        "日志审计：多类型解密查询、条件筛选、手机号匹配、实时查询、CSV导出",
        "数据同步：全量同步、增量同步、定时自动调度",
        "通讯录管理：组织架构树、成员查询、全量/增量同步",
        "密钥管理：多版本 RSA 密钥、在线热切换、密钥测试",
        "操作审计：自动记录所有 API 操作"
    ])

    # 9. 总览看板
    add_content_slide(prs, "功能模块 - 总览看板", [
        "KPI 概览：最新同步时间、近7日同步记录、失败类型数、活跃密钥、通讯录人数、未使用率",
        "最近同步任务：时间、类型、触发方式、成功/失败数、耗时",
        "问题提醒：同步失败、密钥超期、通讯录过期、无新数据",
        "使用分析：按周/月/季统计，部门筛选，未使用人员分析",
        "未使用人员导出：支持按部门、时间范围筛选，CSV导出"
    ])

    # 10. 日志审计
    add_content_slide(prs, "功能模块 - 日志审计查询", [
        "日志类型：下拉多选，支持按数据类型过滤",
        "时间范围：快捷按钮（今天/昨天/近7天/近30天/自定义）",
        "手机号匹配：输入手机号精确查找人员日志",
        "实时查询：开关控制，开启时直接从政务微信实时拉取",
        "自定义条件：支持多组字段+操作符+值筛选",
        "结果展示：按时间倒序，发送方格式「手机号(姓名)」，展开查看完整 JSON"
    ])

    # 11. 数据同步
    add_content_slide(prs, "功能模块 - 数据同步", [
        "定期同步（自动调度）：启动/停止定时增量同步，间隔可配置",
        "全量同步：按时间范围重新拉取，适合数据修复",
        "增量同步：自动从断点继续，只拉取新数据，不丢失",
        "同步状态：进度条 + 当前类型 + 结果表格实时更新",
        "同步历史：追溯记录，排查问题",
        "企微操作日志同步：独立配置，支持设置时间范围"
    ])

    # 12. 通讯录管理
    add_content_slide(prs, "功能模块 - 通讯录管理", [
        "组织架构树：左侧部门树形结构，点击过滤成员",
        "成员列表：姓名、手机号、部门、职位，支持搜索",
        "全量同步：拉取所有部门和成员，首次部署时使用",
        "增量同步：仅拉取新增/变更成员，日常使用",
        "异步导出：后台任务，支持导出通讯录数据",
        "成员详情：点击查看详情，支持查询日志跳转"
    ])

    # 13. 密钥管理
    add_content_slide(prs, "功能模块 - 密钥管理", [
        "密钥列表：展示所有 RSA 密钥版本及状态（活跃/备用）",
        "添加密钥：输入版本号 + PEM 内容，上传新的 RSA 私钥",
        "切换活跃：将备用密钥设为活跃，后续解密优先使用",
        "测试密钥：验证密钥格式是否正确（位数、类型）",
        "密钥热切换：解密失败时自动回退到历史版本密钥",
        "安全机制：PEM 文件权限自动检查修正（600）"
    ])

    # 14. 数据类型配置
    add_content_slide(prs, "功能模块 - 数据类型配置", [
        "功能：管理需要同步的日志数据类型",
        "启用/禁用：每行提供开关，支持批量操作",
        "同步状态：展示各类型的最近同步时间、累计记录数",
        "使用场景：按需启用特定类型，节省资源",
        "点击查看最近日志：直接跳转到该类型的最新记录"
    ])

    # 15. 操作审计
    add_content_slide(prs, "功能模块 - 操作审计", [
        "自动记录：所有 API 操作自动记录，无需手动配置",
        "筛选条件：按操作类型（登录/同步/查询/密钥管理等）和状态过滤",
        "日志字段：操作时间、操作用户、操作类型、请求方法、路径、状态码、耗时",
        "用途：追溯操作、排查异常、满足审计合规要求"
    ])

    # 16. 系统状态
    add_content_slide(prs, "功能模块 - 系统状态", [
        "数据库连通性",
        "系统运行时间",
        "密钥状态（活跃版本 + 总数）",
        "通讯录状态（最近同步时间 + 成员总数）",
        "各数据类型同步覆盖率",
        "数据库存储空间（各表记录数、数据大小、索引大小）"
    ])

    # 17. 安全设计
    add_content_slide(prs, "安全与运维 - 安全设计", [
        "认证机制：JWT HS256 签名，24小时有效期",
        "密码存储：bcrypt 哈希，不可逆",
        "登录限流：1分钟内5次失败自动锁定1分钟",
        "密钥文件安全：启动自动检查权限，宽松权限自动修正为600",
        "可选 PEM 加密：KEY_ENCRYPT_KEY 环境变量启用 AES-256-GCM",
        "操作审计：异步写入，不阻塞请求"
    ])

    # 18. 部署架构
    add_content_slide(prs, "部署架构", [
        "Docker Compose 三容器编排",
        "MySQL：内部端口3306，映射到宿主机3307",
        "后端 Go Echo：内部端口8080，映射到宿主机3010",
        "前端 Nginx：内部端口80，映射到宿主机5173",
        "后端 depends_on MySQL（condition: service_healthy）",
        "前端 depends_on 后端"
    ])

    # 19. 快速开始
    add_content_slide(prs, "快速开始", [
        "1. git clone 源码",
        "2. cp .env.example .env，编辑配置",
        "3. docker-compose up -d --build",
        "4. 验证健康状态：curl http://localhost:3010/health",
        "5. 上传 RSA 密钥并激活",
        "6. 同步通讯录（首次全量同步）",
        "7. 配置数据类型并执行首次同步",
        "8. 启动定时调度"
    ])

    # 20. 日常运维
    add_content_slide(prs, "日常运维", [
        "每日：检查看板异常告警，确认同步调度正常运行",
        "每周：查看同步历史，确认无失败记录，检查存储空间",
        "每月：审查操作日志，评估是否需要更换密钥，检查未使用人员",
        "故障排查：同步卡死→取消重试；解密失败→检查密钥；查询无结果→确认数据类型"
    ])

    # 21. 核心流程 - 解密管道
    add_content_slide(prs, "核心技术 - 解密管道", [
        "Step 1: RSA 解密 encrypt_key (base64) → 16字节 AES 密钥",
        "Step 2: AES-128-CBC 解密 encrypt_data，IV = 密文前16字节",
        "Step 3: PKCS7 去填充，截去尾部8字节",
        "Step 4: JSON 解析得到 LogEntry",
        "密钥回退：解密失败时自动尝试历史版本密钥"
    ])

    # 22. 核心流程 - 数据去重
    add_content_slide(prs, "核心技术 - 数据去重与分表", [
        "去重策略：enc_data_hash (MD5) + INSERT IGNORE",
        "分表策略：log_{feature_id}_{YYYYMM}",
        "动态建表：按 feature ID 和月份自动创建",
        "唯一索引：uk_dedup (feature_id, log_time, enc_data_hash)",
        "虚拟列：login_openid 从 parsed_json 提取，用于手机号快速匹配"
    ])

    # 23. 核心技术 - 增量同步
    add_content_slide(prs, "核心技术 - 增量同步", [
        "记录断点：sync_state.last_log_time",
        "下次同步从断点+1开始，只拉取新数据",
        "即使中间有数据未同步完成，下次仍会从断点继续",
        "定时调度：goroutine + ticker，按配置间隔自动执行",
        "优雅关闭：停止调度 → 取消同步 → 关闭 HTTP → 关闭 DB"
    ])

    # 24. Q&A
    add_title_slide(prs, "Q&A", "常见问题解答\n\n• 同步周期多久？定时调度默认1小时增量同步\n• 密钥更换影响历史数据？不会，仅影响后续新数据\n• 如何知道哪些人没有使用？看板→未使用人员分析\n• 数据会重复存储？不会，使用 MD5 去重\n\n\n感谢使用！")

    # 保存
    output_path = "/Users/leohang/project/wwlocal-wework/docs/政务微信数据审计平台_产品演示.pptx"
    prs.save(output_path)
    print(f"PPTX 已生成: {output_path}")

if __name__ == "__main__":
    main()
